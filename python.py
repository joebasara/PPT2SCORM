import html
import io
import json
import re
import uuid
import zipfile
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

APP_TITLE = "PPTX to SCORM Publisher"


# ============================================================
# General helpers
# ============================================================

def safe_name(value: str) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    return value.strip("._-") or "course"


def emu_pct(value: float, total: float) -> float:
    if not total:
        return 0.0
    return (value / total) * 100.0


def shape_bounds(shape, offset_x: int = 0, offset_y: int = 0) -> Tuple[int, int, int, int]:
    return (
        int(getattr(shape, "left", 0)) + offset_x,
        int(getattr(shape, "top", 0)) + offset_y,
        int(getattr(shape, "width", 0)),
        int(getattr(shape, "height", 0)),
    )


def color_to_hex(color_obj) -> Optional[str]:
    try:
        rgb = getattr(color_obj, "rgb", None)
        if rgb:
            return f"#{rgb}"
    except Exception:
        pass
    return None


def get_auto_shape_type(shape):
    try:
        return shape.auto_shape_type
    except Exception:
        return None


def shape_fill_color(shape) -> str:
    try:
        fill = shape.fill
        fc = fill.fore_color
        color = color_to_hex(fc)
        if color:
            return color
    except Exception:
        pass
    return "transparent"


def shape_line_color(shape) -> str:
    try:
        line = shape.line
        color = color_to_hex(line.color)
        if color:
            return color
    except Exception:
        pass
    return "#000000"


def shape_line_width_px(shape) -> float:
    try:
        width_emu = float(shape.line.width)
        px = width_emu / 12700.0
        return max(px, 1.0)
    except Exception:
        return 1.5


def paragraph_align_to_css(paragraph) -> str:
    try:
        align = paragraph.alignment
        if align == PP_ALIGN.CENTER:
            return "center"
        if align == PP_ALIGN.RIGHT:
            return "right"
        if align == PP_ALIGN.JUSTIFY:
            return "justify"
    except Exception:
        pass
    return "left"


def vertical_anchor_to_css(text_frame, default_center: bool = False) -> str:
    try:
        va = text_frame.vertical_anchor
        if va == MSO_ANCHOR.TOP:
            return "flex-start"
        if va == MSO_ANCHOR.MIDDLE:
            return "center"
        if va == MSO_ANCHOR.BOTTOM:
            return "flex-end"
    except Exception:
        pass
    return "center" if default_center else "flex-start"


# ============================================================
# Hyperlink detection
# ============================================================

def extract_shape_external_link(shape) -> Optional[str]:
    try:
        click_action = getattr(shape, "click_action", None)
        if click_action is not None:
            hyperlink = getattr(click_action, "hyperlink", None)
            if hyperlink is not None:
                address = getattr(hyperlink, "address", None)
                if address:
                    return address
    except Exception:
        pass
    return None


def detect_internal_link_target(slides, slide_idx_zero: int, shape) -> Optional[int]:
    try:
        click_action = getattr(shape, "click_action", None)
        if click_action is not None:
            target_slide = getattr(click_action, "target_slide", None)
            if target_slide is not None:
                for i, s in enumerate(slides, start=1):
                    if s == target_slide:
                        return i
    except Exception:
        pass

    try:
        hlink_click = shape.element.xpath(".//*[local-name()='hlinkClick']")
        if hlink_click:
            action = (hlink_click[0].get("action") or "").lower()
            if "firstslide" in action:
                return 1
            if "lastslide" in action:
                return len(slides)
            if "nextslide" in action:
                return min(len(slides), slide_idx_zero + 2)
            if "previousslide" in action:
                return max(1, slide_idx_zero)
    except Exception:
        pass

    return None


# ============================================================
# Shape-aware hotspot geometry
# ============================================================

def arrow_polygon_points(auto, x, y, w, h) -> Optional[List[Tuple[float, float]]]:
    if auto == MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW:
        return [
            (x, y + h * 0.25),
            (x + w * 0.65, y + h * 0.25),
            (x + w * 0.65, y),
            (x + w, y + h * 0.5),
            (x + w * 0.65, y + h),
            (x + w * 0.65, y + h * 0.75),
            (x, y + h * 0.75),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.LEFT_ARROW:
        return [
            (x + w, y + h * 0.25),
            (x + w * 0.35, y + h * 0.25),
            (x + w * 0.35, y),
            (x, y + h * 0.5),
            (x + w * 0.35, y + h),
            (x + w * 0.35, y + h * 0.75),
            (x + w, y + h * 0.75),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.UP_ARROW:
        return [
            (x + w * 0.25, y + h),
            (x + w * 0.25, y + h * 0.35),
            (x, y + h * 0.35),
            (x + w * 0.5, y),
            (x + w, y + h * 0.35),
            (x + w * 0.75, y + h * 0.35),
            (x + w * 0.75, y + h),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.DOWN_ARROW:
        return [
            (x + w * 0.25, y),
            (x + w * 0.25, y + h * 0.65),
            (x, y + h * 0.65),
            (x + w * 0.5, y + h),
            (x + w, y + h * 0.65),
            (x + w * 0.75, y + h * 0.65),
            (x + w * 0.75, y),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.LEFT_RIGHT_ARROW:
        return [
            (x + w * 0.18, y + h * 0.25),
            (x + w * 0.82, y + h * 0.25),
            (x + w * 0.82, y),
            (x + w, y + h * 0.5),
            (x + w * 0.82, y + h),
            (x + w * 0.82, y + h * 0.75),
            (x + w * 0.18, y + h * 0.75),
            (x + w * 0.18, y + h),
            (x, y + h * 0.5),
            (x + w * 0.18, y),
        ]
    return None


def regular_polygon_points(sides: int, x: float, y: float, w: float, h: float, rotate_deg: float = -90.0):
    import math
    cx = x + w / 2
    cy = y + h / 2
    rx = w / 2
    ry = h / 2
    pts = []
    for i in range(sides):
        ang = math.radians(rotate_deg + (360 / sides) * i)
        pts.append((cx + rx * math.cos(ang), cy + ry * math.sin(ang)))
    return pts


def custom_polygon_points(auto, x, y, w, h) -> Optional[List[Tuple[float, float]]]:
    if auto == MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE:
        return [(x + w / 2, y), (x + w, y + h), (x, y + h)]
    if auto == MSO_AUTO_SHAPE_TYPE.DIAMOND:
        return [(x + w / 2, y), (x + w, y + h / 2), (x + w / 2, y + h), (x, y + h / 2)]
    if auto == MSO_AUTO_SHAPE_TYPE.PENTAGON:
        return regular_polygon_points(5, x, y, w, h)
    if auto == MSO_AUTO_SHAPE_TYPE.HEXAGON:
        return [
            (x + w * 0.25, y),
            (x + w * 0.75, y),
            (x + w, y + h * 0.5),
            (x + w * 0.75, y + h),
            (x + w * 0.25, y + h),
            (x, y + h * 0.5),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.CHEVRON:
        return [
            (x, y),
            (x + w * 0.65, y),
            (x + w, y + h * 0.5),
            (x + w * 0.65, y + h),
            (x, y + h),
            (x + w * 0.35, y + h * 0.5),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.PARALLELOGRAM:
        return [
            (x + w * 0.2, y),
            (x + w, y),
            (x + w * 0.8, y + h),
            (x, y + h),
        ]
    if auto == MSO_AUTO_SHAPE_TYPE.TRAPEZOID:
        return [
            (x + w * 0.2, y),
            (x + w * 0.8, y),
            (x + w, y + h),
            (x, y + h),
        ]
    return None


def hotspot_geometry_for_shape(shape, offset_x=0, offset_y=0) -> Dict[str, Any]:
    x, y, w, h = shape_bounds(shape, offset_x, offset_y)
    auto = get_auto_shape_type(shape)
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.LINE:
        return {
            "geom": "line",
            "x1": x,
            "y1": y,
            "x2": x + w,
            "y2": y + h,
            "x": x,
            "y": y,
            "w": w,
            "h": h,
        }

    if auto == MSO_AUTO_SHAPE_TYPE.OVAL:
        return {"geom": "ellipse", "x": x, "y": y, "w": w, "h": h}

    if auto == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
        return {"geom": "rect", "x": x, "y": y, "w": w, "h": h}

    if auto == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
        return {
            "geom": "roundrect",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "rx": min(w, h) * 0.12,
            "ry": min(w, h) * 0.12,
        }

    pts = arrow_polygon_points(auto, x, y, w, h)
    if pts:
        return {"geom": "polygon", "x": x, "y": y, "w": w, "h": h, "points": pts}

    pts = custom_polygon_points(auto, x, y, w, h)
    if pts:
        return {"geom": "polygon", "x": x, "y": y, "w": w, "h": h, "points": pts}

    return {"geom": "rect", "x": x, "y": y, "w": w, "h": h}


# ============================================================
# SVG shape rendering
# ============================================================

def svg_shape_spec(shape, slide_w, slide_h, offset_x=0, offset_y=0) -> Optional[Dict[str, Any]]:
    x, y, w, h = shape_bounds(shape, offset_x, offset_y)
    fill = shape_fill_color(shape)
    stroke = shape_line_color(shape)
    stroke_width = shape_line_width_px(shape)
    auto = get_auto_shape_type(shape)
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.LINE:
        return {
            "kind": "line",
            "x1": x,
            "y1": y,
            "x2": x + w,
            "y2": y + h,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    if auto == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
        return {
            "kind": "rect",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "rx": 0,
            "ry": 0,
            "fill": fill,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    if auto == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
        return {
            "kind": "rect",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "rx": min(w, h) * 0.12,
            "ry": min(w, h) * 0.12,
            "fill": fill,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    if auto == MSO_AUTO_SHAPE_TYPE.OVAL:
        return {
            "kind": "ellipse",
            "cx": x + w / 2,
            "cy": y + h / 2,
            "rx": w / 2,
            "ry": h / 2,
            "fill": fill,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    pts = arrow_polygon_points(auto, x, y, w, h)
    if pts:
        return {
            "kind": "polygon",
            "points": pts,
            "fill": fill,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    pts = custom_polygon_points(auto, x, y, w, h)
    if pts:
        return {
            "kind": "polygon",
            "points": pts,
            "fill": fill,
            "stroke": stroke,
            "stroke_width": stroke_width,
        }

    return None


# ============================================================
# Text extraction
# ============================================================

def run_to_html(run) -> str:
    txt = html.escape(run.text or "")
    if not txt:
        return ""

    styles = []
    try:
        font = run.font
        if getattr(font, "bold", False):
            styles.append("font-weight:700")
        if getattr(font, "italic", False):
            styles.append("font-style:italic")
        if getattr(font, "underline", False):
            styles.append("text-decoration:underline")
        size = getattr(font, "size", None)
        if size is not None:
            styles.append(f"font-size:{max(10, int(size.pt))}px")
        try:
            color = color_to_hex(font.color)
            if color:
                styles.append(f"color:{color}")
        except Exception:
            pass
    except Exception:
        pass

    if styles:
        txt = f'<span style="{";".join(styles)}">{txt}</span>'

    try:
        url = run.hyperlink.address
    except Exception:
        url = None

    if url:
        txt = f'<a href="{html.escape(url)}" target="_blank" rel="noopener noreferrer">{txt}</a>'

    return txt


def extract_text_element(shape, offset_x=0, offset_y=0, default_center=False) -> Optional[Dict[str, Any]]:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return None

    x, y, w, h = shape_bounds(shape, offset_x, offset_y)
    paragraphs = []

    for para in text_frame.paragraphs:
        parts = [run_to_html(run) for run in para.runs]
        parts = [p for p in parts if p]
        if parts:
            paragraphs.append({
                "align": paragraph_align_to_css(para),
                "html": "".join(parts),
            })

    if not paragraphs:
        return None

    return {
        "type": "text",
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "paragraphs": paragraphs,
        "v_align": vertical_anchor_to_css(text_frame, default_center=default_center),
    }


# ============================================================
# Table extraction
# ============================================================

def extract_table_element(shape, offset_x=0, offset_y=0) -> Optional[Dict[str, Any]]:
    try:
        table = shape.table
    except Exception:
        return None

    x, y, w, h = shape_bounds(shape, offset_x, offset_y)

    try:
        col_widths = [int(col.width) for col in table.columns]
    except Exception:
        col_widths = []

    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_html_parts = []
            for para in cell.text_frame.paragraphs:
                parts = [run_to_html(run) for run in para.runs]
                parts = [p for p in parts if p]
                if parts:
                    cell_html_parts.append(
                        f'<div style="text-align:{paragraph_align_to_css(para)};">{"".join(parts)}</div>'
                    )

            fill = "transparent"
            border = "#555"

            try:
                fc = color_to_hex(cell.fill.fore_color)
                if fc:
                    fill = fc
            except Exception:
                pass

            cells.append({
                "html": "".join(cell_html_parts) if cell_html_parts else "&nbsp;",
                "fill": fill,
                "border": border,
            })
        rows.append(cells)

    return {
        "type": "table",
        "x": x,
        "y": y,
        "w": w,
        "h": h,
        "col_widths": col_widths,
        "rows": rows,
    }


# ============================================================
# Image extraction
# ============================================================

def extract_image_element(shape, media: Dict[str, bytes], slide_index: int, image_index: int, offset_x=0, offset_y=0):
    try:
        image = getattr(shape, "image", None)
        if image is None:
            return None
        ext = image.ext or "png"
        filename = f"media/slide_{slide_index:03d}_img_{image_index:02d}.{ext}"
        media[filename] = image.blob
        x, y, w, h = shape_bounds(shape, offset_x, offset_y)
        return {
            "type": "image",
            "x": x,
            "y": y,
            "w": w,
            "h": h,
            "src": filename,
        }
    except Exception:
        return None


# ============================================================
# Recursive slide extraction
# ============================================================

def process_shape(
    shape,
    prs: Presentation,
    slide_idx_zero: int,
    slide_out: Dict[str, Any],
    media: Dict[str, bytes],
    slide_index_one: int,
    counters: Dict[str, int],
    offset_x: int = 0,
    offset_y: int = 0,
):
    shape_type = getattr(shape, "shape_type", None)

    if shape_type == MSO_SHAPE_TYPE.GROUP:
        gx, gy, _, _ = shape_bounds(shape, offset_x, offset_y)
        try:
            for child in shape.shapes:
                process_shape(
                    child,
                    prs,
                    slide_idx_zero,
                    slide_out,
                    media,
                    slide_index_one,
                    counters,
                    offset_x=gx,
                    offset_y=gy,
                )
        except Exception:
            pass
        return

    img_el = extract_image_element(shape, media, slide_index_one, counters["images"] + 1, offset_x, offset_y)
    if img_el is not None:
        counters["images"] += 1
        slide_out["images"].append(img_el)

    table_el = extract_table_element(shape, offset_x, offset_y)
    if table_el is not None:
        slide_out["tables"].append(table_el)

    svg = svg_shape_spec(shape, int(prs.slide_width), int(prs.slide_height), offset_x, offset_y)
    if svg is not None:
        slide_out["svg_shapes"].append(svg)

    default_center = get_auto_shape_type(shape) is not None
    text_el = extract_text_element(shape, offset_x, offset_y, default_center=default_center)
    if text_el is not None:
        slide_out["text_elements"].append(text_el)

    try:
        external = extract_shape_external_link(shape)
        internal = detect_internal_link_target(prs.slides, slide_idx_zero, shape)
        if external or internal:
            geom = hotspot_geometry_for_shape(shape, offset_x, offset_y)
            if external:
                geom = {**geom, "kind": "external", "url": external}
            else:
                geom = {**geom, "kind": "internal", "target_slide": internal}
            slide_out["hotspots"].append(geom)
    except Exception:
        pass


def extract_course(prs: Presentation) -> Tuple[Dict[str, Any], Dict[str, bytes]]:
    media: Dict[str, bytes] = {}
    slides_out = []

    for s_idx, slide in enumerate(prs.slides, start=1):
        slide_out = {
            "index": s_idx,
            "svg_shapes": [],
            "images": [],
            "tables": [],
            "text_elements": [],
            "hotspots": [],
        }
        counters = {"images": 0}

        for shape in slide.shapes:
            process_shape(
                shape=shape,
                prs=prs,
                slide_idx_zero=s_idx - 1,
                slide_out=slide_out,
                media=media,
                slide_index_one=s_idx,
                counters=counters,
                offset_x=0,
                offset_y=0,
            )

        slides_out.append(slide_out)

    course = {
        "slideWidthEmu": int(prs.slide_width),
        "slideHeightEmu": int(prs.slide_height),
        "slides": slides_out,
    }
    return course, media


# ============================================================
# SCORM
# ============================================================

def build_scorm_driver_js() -> str:
    return """
var scormAPI = null;

function findAPI(win) {
  var tries = 0;
  while ((win.API == null) && (win.parent != null) && (win.parent != win)) {
    tries++;
    if (tries > 20) return null;
    win = win.parent;
  }
  return win.API;
}

function getAPI() {
  if (scormAPI == null) scormAPI = findAPI(window);
  return scormAPI;
}

function scormInitialize() {
  var api = getAPI();
  if (api) return api.LMSInitialize("");
  return false;
}

function scormTerminate() {
  var api = getAPI();
  if (api) return api.LMSFinish("");
  return false;
}

function scormGetValue(key) {
  var api = getAPI();
  if (api) return api.LMSGetValue(key);
  return "";
}

function scormSetValue(key, value) {
  var api = getAPI();
  if (api) return api.LMSSetValue(key, value);
  return false;
}

function scormCommit() {
  var api = getAPI();
  if (api) return api.LMSCommit("");
  return false;
}
"""


def build_manifest_xml(course_id: str, title: str, media_files: List[str]) -> str:
    media_xml = "\n".join([f'      <file href="{name}" />' for name in media_files])

    return f"""<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="{course_id}"
    version="1.0"
    xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
    xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2"
    xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance"
    xsi:schemaLocation="http://www.imsproject.org/xsd/imscp_rootv1p1p2 imscp_rootv1p1p2.xsd
    http://www.adlnet.org/xsd/adlcp_rootv1p2 adlcp_rootv1p2.xsd">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="ORG1">
    <organization identifier="ORG1">
      <title>{html.escape(title)}</title>
      <item identifier="ITEM1" identifierref="RES1">
        <title>{html.escape(title)}</title>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="RES1" type="webcontent" adlcp:scormtype="sco" href="index_lms.html">
      <file href="index_lms.html" />
      <file href="scormdriver.js" />
{media_xml}
    </resource>
  </resources>
</manifest>"""


# ============================================================
# HTML player
# ============================================================

def build_player_html(title: str, course: Dict[str, Any]) -> str:
    data = json.dumps(course, ensure_ascii=False)

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>{html.escape(title)}</title>
<style>
* {{
  box-sizing: border-box;
}}

body {{
  margin: 0;
  font-family: Arial, sans-serif;
  background: #111;
  color: #f1f1f1;
  min-height: 100vh;
}}

.app {{
  display: flex;
  min-height: 100vh;
}}

.sidebar {{
  width: 130px;
  background: #1a1a1a;
  border-right: 1px solid #333;
  padding: 16px 10px;
}}

.sidebar label {{
  display: block;
  font-size: 13px;
  color: #bbb;
  margin-bottom: 8px;
}}

.sidebar select {{
  width: 100%;
  height: calc(100vh - 40px);
  min-height: 260px;
  background: #262626;
  color: #f1f1f1;
  border: 1px solid #444;
  border-radius: 8px;
  padding: 8px;
}}

.main {{
  flex: 1;
  display: flex;
  flex-direction: column;
}}

.header {{
  background: #1a1a1a;
  border-bottom: 1px solid #333;
  padding: 12px 16px;
}}

.header-title {{
  font-weight: 700;
}}

.header-sub {{
  color: #bbb;
  font-size: 14px;
  margin-top: 4px;
}}

.stage-wrap {{
  flex: 1;
  display: grid;
  place-items: center;
  padding: 18px;
}}

.frame {{
  width: min(100%, 1280px);
  aspect-ratio: 16 / 9;
  position: relative;
  background: white;
  overflow: hidden;
  border-radius: 10px;
  box-shadow: 0 14px 40px rgba(0,0,0,.45);
}}

.svg-layer,
.image-layer,
.table-layer,
.text-layer,
.hotspot-layer {{
  position: absolute;
  inset: 0;
}}

.svg-layer svg {{
  width: 100%;
  height: 100%;
  display: block;
}}

.el-img {{
  position: absolute;
  object-fit: contain;
}}

.el-text {{
  position: absolute;
  display: flex;
  flex-direction: column;
  overflow: hidden;
  color: #111;
  line-height: 1.2;
  pointer-events: auto;
  padding: 2px 4px;
  white-space: normal;
  word-break: break-word;
}}

.el-text a {{
  color: #0a58ca;
  text-decoration: underline;
}}

.el-table {{
  position: absolute;
  border-collapse: collapse;
  table-layout: fixed;
  background: transparent;
  color: #111;
  font-size: 14px;
}}

.el-table td {{
  border: 1px solid #555;
  padding: 4px 6px;
  vertical-align: middle;
  overflow: hidden;
}}

.hotspot-layer {{
  pointer-events: none;
}}

.hotspot-shape {{
  position: absolute;
  inset: 0;
  width: 100%;
  height: 100%;
  pointer-events: auto;
}}

.hotspot-shape a {{
  cursor: pointer;
}}

.bottom-bar {{
  background: #1a1a1a;
  border-top: 1px solid #333;
  padding: 12px 16px;
  display: flex;
  justify-content: center;
  gap: 12px;
}}

.bottom-bar button {{
  background: #262626;
  color: #f1f1f1;
  border: 1px solid #444;
  border-radius: 8px;
  padding: 10px 16px;
  cursor: pointer;
  min-width: 100px;
}}

.footer {{
  text-align: center;
  background: #1a1a1a;
  color: #bbb;
  font-size: 13px;
  padding: 8px 16px 14px;
}}
</style>
</head>
<body>
<div class="app">
  <aside class="sidebar">
    <label for="jumpSelect">Slides</label>
    <select id="jumpSelect" size="20" onchange="jumpSlide(this.value)"></select>
  </aside>

  <div class="main">
    <div class="header">
      <div class="header-title">{html.escape(title)}</div>
      <div class="header-sub" id="counter">Slide 1</div>
    </div>

    <div class="stage-wrap">
      <div class="frame">
        <div class="svg-layer" id="svgLayer"></div>
        <div class="image-layer" id="imageLayer"></div>
        <div class="table-layer" id="tableLayer"></div>
        <div class="text-layer" id="textLayer"></div>
        <div class="hotspot-layer" id="hotspotLayer"></div>
      </div>
    </div>

    <div class="bottom-bar">
      <button onclick="prevSlide()">Previous</button>
      <button onclick="nextSlide()">Next</button>
    </div>

    <div class="footer">Pure-Python SVG/HTML renderer with shape-aware hotspots.</div>
  </div>
</div>

<script src="scormdriver.js"></script>
<script>
const course = {data};
let currentSlide = 1;

const svgLayer = document.getElementById("svgLayer");
const imageLayer = document.getElementById("imageLayer");
const tableLayer = document.getElementById("tableLayer");
const textLayer = document.getElementById("textLayer");
const hotspotLayer = document.getElementById("hotspotLayer");
const counter = document.getElementById("counter");
const jumpSelect = document.getElementById("jumpSelect");

function pct(value, total) {{
  return (value / total) * 100;
}}

function populateJumpMenu() {{
  jumpSelect.innerHTML = "";
  for (const slide of course.slides) {{
    const opt = document.createElement("option");
    opt.value = slide.index;
    opt.textContent = `Slide ${{slide.index}}`;
    jumpSelect.appendChild(opt);
  }}
}}

function renderSVG(slide) {{
  let shapes = "";
  for (const raw of slide.svg_shapes || []) {{
    if (raw.kind === "rect") {{
      shapes += `<rect x="${{pct(raw.x, course.slideWidthEmu)}}%" y="${{pct(raw.y, course.slideHeightEmu)}}%" width="${{pct(raw.w, course.slideWidthEmu)}}%" height="${{pct(raw.h, course.slideHeightEmu)}}%" rx="${{pct(raw.rx || 0, course.slideWidthEmu)}}%" ry="${{pct(raw.ry || 0, course.slideHeightEmu)}}%" fill="${{raw.fill}}" stroke="${{raw.stroke}}" stroke-width="${{raw.stroke_width}}" />`;
    }} else if (raw.kind === "ellipse") {{
      shapes += `<ellipse cx="${{pct(raw.cx, course.slideWidthEmu)}}%" cy="${{pct(raw.cy, course.slideHeightEmu)}}%" rx="${{pct(raw.rx, course.slideWidthEmu)}}%" ry="${{pct(raw.ry, course.slideHeightEmu)}}%" fill="${{raw.fill}}" stroke="${{raw.stroke}}" stroke-width="${{raw.stroke_width}}" />`;
    }} else if (raw.kind === "line") {{
      shapes += `<line x1="${{pct(raw.x1, course.slideWidthEmu)}}%" y1="${{pct(raw.y1, course.slideHeightEmu)}}%" x2="${{pct(raw.x2, course.slideWidthEmu)}}%" y2="${{pct(raw.y2, course.slideHeightEmu)}}%" stroke="${{raw.stroke}}" stroke-width="${{raw.stroke_width}}" />`;
    }} else if (raw.kind === "polygon") {{
      const pts = (raw.points || []).map(p => `${{pct(p[0], course.slideWidthEmu)}},${{pct(p[1], course.slideHeightEmu)}}`).join(" ");
      shapes += `<polygon points="${{pts}}" fill="${{raw.fill}}" stroke="${{raw.stroke}}" stroke-width="${{raw.stroke_width}}" />`;
    }}
  }}

  svgLayer.innerHTML = `<svg viewBox="0 0 100 100" preserveAspectRatio="none">${{shapes}}</svg>`;
}}

function renderImages(slide) {{
  imageLayer.innerHTML = "";
  for (const el of slide.images || []) {{
    const img = document.createElement("img");
    img.className = "el-img";
    img.src = el.src;
    img.style.left = pct(el.x, course.slideWidthEmu) + "%";
    img.style.top = pct(el.y, course.slideHeightEmu) + "%";
    img.style.width = pct(el.w, course.slideWidthEmu) + "%";
    img.style.height = pct(el.h, course.slideHeightEmu) + "%";
    imageLayer.appendChild(img);
  }}
}}

function renderTables(slide) {{
  tableLayer.innerHTML = "";
  for (const t of slide.tables || []) {{
    const table = document.createElement("table");
    table.className = "el-table";
    table.style.left = pct(t.x, course.slideWidthEmu) + "%";
    table.style.top = pct(t.y, course.slideHeightEmu) + "%";
    table.style.width = pct(t.w, course.slideWidthEmu) + "%";
    table.style.height = pct(t.h, course.slideHeightEmu) + "%";

    if (t.col_widths && t.col_widths.length) {{
      const cg = document.createElement("colgroup");
      const total = t.col_widths.reduce((a,b)=>a+b,0);
      for (const cw of t.col_widths) {{
        const col = document.createElement("col");
        col.style.width = pct(cw, total) + "%";
        cg.appendChild(col);
      }}
      table.appendChild(cg);
    }}

    for (const row of t.rows || []) {{
      const tr = document.createElement("tr");
      for (const cell of row) {{
        const td = document.createElement("td");
        td.innerHTML = cell.html || "&nbsp;";
        td.style.background = cell.fill || "transparent";
        td.style.borderColor = cell.border || "#555";
        tr.appendChild(td);
      }}
      table.appendChild(tr);
    }}

    tableLayer.appendChild(table);
  }}
}}

function renderText(slide) {{
  textLayer.innerHTML = "";
  for (const el of slide.text_elements || []) {{
    const div = document.createElement("div");
    div.className = "el-text";
    div.style.left = pct(el.x, course.slideWidthEmu) + "%";
    div.style.top = pct(el.y, course.slideHeightEmu) + "%";
    div.style.width = pct(el.w, course.slideWidthEmu) + "%";
    div.style.height = pct(el.h, course.slideHeightEmu) + "%";
    div.style.justifyContent = el.v_align || "flex-start";

    let inner = "";
    for (const p of el.paragraphs || []) {{
      inner += `<div style="text-align:${{p.align || "left"}};">${{p.html}}</div>`;
    }}
    div.innerHTML = inner;
    textLayer.appendChild(div);
  }}
}}

function addHotspotAnchor(el, link) {{
  if (link.kind === "external" && link.url) {{
    el.setAttribute("href", link.url);
    el.setAttribute("target", "_blank");
    el.setAttribute("rel", "noopener noreferrer");
  }} else if (link.kind === "internal" && link.target_slide) {{
    el.setAttribute("href", "#");
    el.addEventListener("click", function(e) {{
      e.preventDefault();
      goToSlide(Number(link.target_slide));
    }});
  }}
  return el;
}}

function renderHotspots(slide) {{
  hotspotLayer.innerHTML = "";

  const svg = document.createElementNS("http://www.w3.org/2000/svg", "svg");
  svg.setAttribute("class", "hotspot-shape");
  svg.setAttribute("viewBox", "0 0 100 100");
  svg.setAttribute("preserveAspectRatio", "none");

  for (const link of slide.hotspots || []) {{
    let anchor = document.createElementNS("http://www.w3.org/2000/svg", "a");
    addHotspotAnchor(anchor, link);

    let shapeEl = null;

    if (link.geom === "ellipse") {{
      shapeEl = document.createElementNS("http://www.w3.org/2000/svg", "ellipse");
      shapeEl.setAttribute("cx", pct(link.x + link.w / 2, course.slideWidthEmu));
      shapeEl.setAttribute("cy", pct(link.y + link.h / 2, course.slideHeightEmu));
      shapeEl.setAttribute("rx", pct(link.w / 2, course.slideWidthEmu));
      shapeEl.setAttribute("ry", pct(link.h / 2, course.slideHeightEmu));
    }} else if (link.geom === "roundrect" || link.geom === "rect") {{
      shapeEl = document.createElementNS("http://www.w3.org/2000/svg", "rect");
      shapeEl.setAttribute("x", pct(link.x, course.slideWidthEmu));
      shapeEl.setAttribute("y", pct(link.y, course.slideHeightEmu));
      shapeEl.setAttribute("width", pct(link.w, course.slideWidthEmu));
      shapeEl.setAttribute("height", pct(link.h, course.slideHeightEmu));
      if (link.geom === "roundrect") {{
        shapeEl.setAttribute("rx", pct(link.rx || 0, course.slideWidthEmu));
        shapeEl.setAttribute("ry", pct(link.ry || 0, course.slideHeightEmu));
      }}
    }} else if (link.geom === "polygon" && link.points) {{
      shapeEl = document.createElementNS("http://www.w3.org/2000/svg", "polygon");
      const pts = link.points.map(p => `${{pct(p[0], course.slideWidthEmu)}},${{pct(p[1], course.slideHeightEmu)}}`).join(" ");
      shapeEl.setAttribute("points", pts);
    }} else if (link.geom === "line") {{
      shapeEl = document.createElementNS("http://www.w3.org/2000/svg", "line");
      shapeEl.setAttribute("x1", pct(link.x1, course.slideWidthEmu));
      shapeEl.setAttribute("y1", pct(link.y1, course.slideHeightEmu));
      shapeEl.setAttribute("x2", pct(link.x2, course.slideWidthEmu));
      shapeEl.setAttribute("y2", pct(link.y2, course.slideHeightEmu));
      shapeEl.setAttribute("stroke-width", "12");
      shapeEl.setAttribute("stroke", "rgba(0,0,0,0)");
      shapeEl.setAttribute("pointer-events", "stroke");
      anchor.appendChild(shapeEl);
      svg.appendChild(anchor);
      continue;
    }} else {{
      shapeEl = document.createElementNS("http://www.w3.org/2000/svg", "rect");
      shapeEl.setAttribute("x", pct(link.x, course.slideWidthEmu));
      shapeEl.setAttribute("y", pct(link.y, course.slideHeightEmu));
      shapeEl.setAttribute("width", pct(link.w, course.slideWidthEmu));
      shapeEl.setAttribute("height", pct(link.h, course.slideHeightEmu));
    }}

    shapeEl.setAttribute("fill", "rgba(0,0,0,0)");
    shapeEl.setAttribute("stroke", "rgba(0,0,0,0)");
    anchor.appendChild(shapeEl);
    svg.appendChild(anchor);
  }}

  hotspotLayer.appendChild(svg);
}}

function setScormState() {{
  if (!window.scormSetValue) return;
  try {{
    window.scormSetValue("cmi.core.lesson_location", String(currentSlide));
    const progress = Math.round((currentSlide / course.slides.length) * 100);
    window.scormSetValue("cmi.core.score.raw", String(progress));
    window.scormSetValue("cmi.core.lesson_status", currentSlide >= course.slides.length ? "completed" : "incomplete");
    window.scormCommit();
  }} catch (e) {{}}
}}

function goToSlide(num) {{
  if (num < 1 || num > course.slides.length) return;
  currentSlide = num;
  const slide = course.slides[currentSlide - 1];

  renderSVG(slide);
  renderImages(slide);
  renderTables(slide);
  renderText(slide);
  renderHotspots(slide);

  counter.textContent = `Slide ${{currentSlide}} of ${{course.slides.length}}`;
  jumpSelect.value = String(currentSlide);
  setScormState();
}}

function nextSlide() {{
  if (currentSlide < course.slides.length) goToSlide(currentSlide + 1);
}}

function prevSlide() {{
  if (currentSlide > 1) goToSlide(currentSlide - 1);
}}

function jumpSlide(value) {{
  const num = Number(value);
  if (!Number.isNaN(num)) goToSlide(num);
}}

document.addEventListener("keydown", function(e) {{
  if (e.key === "ArrowRight") nextSlide();
  if (e.key === "ArrowLeft") prevSlide();
}});

window.addEventListener("load", function() {{
  populateJumpMenu();
  if (window.scormInitialize) {{
    try {{
      window.scormInitialize();
      const saved = window.scormGetValue("cmi.core.lesson_location");
      const n = parseInt(saved || "1", 10);
      if (!Number.isNaN(n) && n >= 1 && n <= course.slides.length) currentSlide = n;
    }} catch (e) {{}}
  }}
  goToSlide(currentSlide);
}});

window.addEventListener("beforeunload", function() {{
  if (window.scormTerminate) {{
    try {{ window.scormTerminate(); }} catch (e) {{}}
  }}
}});
</script>
</body>
</html>"""


# ============================================================
# Build ZIP
# ============================================================

def build_scorm_zip(pptx_bytes: bytes, pptx_name: str, course_title: str):
    prs = Presentation(io.BytesIO(pptx_bytes))
    course, media = extract_course(prs)

    mem = io.BytesIO()
    with zipfile.ZipFile(mem, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, blob in media.items():
            zf.writestr(name, blob)

        zf.writestr("index_lms.html", build_player_html(course_title, course))
        zf.writestr("scormdriver.js", build_scorm_driver_js())
        zf.writestr(
            "imsmanifest.xml",
            build_manifest_xml(
                f"{safe_name(course_title)}_{uuid.uuid4().hex[:8]}",
                course_title,
                sorted(media.keys()),
            ),
        )

    mem.seek(0)
    out_name = f"{safe_name(Path(pptx_name).stem)}_scorm12.zip"
    summary = {
        "slides": len(prs.slides),
        "media": len(media),
    }
    return mem.read(), out_name, summary


# ============================================================
# Streamlit UI
# ============================================================

st.set_page_config(page_title=APP_TITLE, layout="centered")
st.title(APP_TITLE)
st.caption("Pure-Python MVP: PPTX → HTML/SVG slide player → SCORM 1.2 ZIP")

with st.expander("What this version supports"):
    st.markdown(
        """
- Basic geometric shape rendering with SVG
- Rectangles, rounded rectangles, ellipses, arrows, several polygon-like shapes
- Text boxes and text inside shapes
- Basic tables
- Embedded pictures
- External hyperlinks
- Internal slide-jump hotspots
- Shape-aware hotspots for supported shapes
- Left slide menu and bottom centered navigation

Still limited:
- Not full PowerPoint fidelity
- No animations, gradients, shadows, or advanced effects
- Complex grouped layouts may still need refinement
- Unsupported shape types fall back to simpler handling
        """
    )

course_title = st.text_input("Course title", value="My SCORM Course")
uploaded = st.file_uploader("Upload PPTX", type=["pptx"])

if st.button("Publish SCORM", type="primary", use_container_width=True):
    if not uploaded:
        st.error("Please upload a .pptx file first.")
    else:
        try:
            zip_bytes, zip_name, summary = build_scorm_zip(
                uploaded.getvalue(),
                uploaded.name,
                course_title,
            )
            st.success("SCORM package created.")
            st.write(f"Slides detected: {summary['slides']}")
            st.write(f"Embedded pictures copied: {summary['media']}")
            st.download_button(
                "Download SCORM ZIP",
                data=zip_bytes,
                file_name=zip_name,
                mime="application/zip",
                use_container_width=True,
            )
        except Exception as e:
            st.error(f"Failed to publish SCORM package: {e}")

st.divider()
st.subheader("requirements.txt")
st.code("streamlit\npython-pptx\nlxml\nPillow\n")
